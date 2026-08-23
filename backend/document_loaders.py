import os
import json
import pandas as pd
from typing import List
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

def extrair_documentos_pdf(caminho_arquivo: str) -> List[Document]:
    carregador = PyPDFLoader(caminho_arquivo)
    return carregador.load()

def extrair_documentos_docx(caminho_arquivo: str) -> List[Document]:
    import docx
    doc = docx.Document(caminho_arquivo)
    texto_completo = []
    for p in doc.paragraphs:
        if p.text.strip():
            texto_completo.append(p.text)
    
    # Extrai tabelas se existirem
    for tabela in doc.tables:
        for linha in tabela.rows:
            texto_linha = " | ".join([celula.text.strip() for celula in linha.cells if celula.text.strip()])
            if texto_linha:
                texto_completo.append(texto_linha)

    conteudo = "\n".join(texto_completo)
    return [Document(page_content=conteudo, metadata={"source": os.path.basename(caminho_arquivo), "format": "Word (.docx)"})]

def extrair_documentos_excel(caminho_arquivo: str) -> List[Document]:
    excel_file = pd.ExcelFile(caminho_arquivo)
    documentos = []
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        texto = f"Aba/Planilha: {sheet_name}\n" + df.to_string(index=False)
        documentos.append(Document(
            page_content=texto,
            metadata={"source": os.path.basename(caminho_arquivo), "sheet": sheet_name, "format": "Excel (.xlsx)"}
        ))
    return documentos

def extrair_documentos_pptx(caminho_arquivo: str) -> List[Document]:
    from pptx import Presentation
    prs = Presentation(caminho_arquivo)
    documentos = []
    for idx, slide in enumerate(prs.slides, start=1):
        texto_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texto_slide.append(shape.text.strip())
        if texto_slide:
            documentos.append(Document(
                page_content="\n".join(texto_slide),
                metadata={"source": os.path.basename(caminho_arquivo), "slide": idx, "format": "PowerPoint (.pptx)"}
            ))
    return documentos

def extrair_documentos_csv(caminho_arquivo: str) -> List[Document]:
    df = pd.read_csv(caminho_arquivo)
    texto = df.to_string(index=False)
    return [Document(page_content=texto, metadata={"source": os.path.basename(caminho_arquivo), "format": "CSV (.csv)"})]

def extrair_documentos_json(caminho_arquivo: str) -> List[Document]:
    with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
        dados = json.load(f)
    texto = json.dumps(dados, indent=2, ensure_ascii=False)
    return [Document(page_content=texto, metadata={"source": os.path.basename(caminho_arquivo), "format": "JSON (.json)"})]

def extrair_documentos_html(caminho_arquivo: str) -> List[Document]:
    from bs4 import BeautifulSoup
    with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    # Remove scripts e estilos
    for s in soup(['script', 'style', 'nav', 'footer']):
        s.decompose()
    texto = soup.get_text(separator='\n')
    linhas = [linha.strip() for linha in texto.splitlines() if linha.strip()]
    return [Document(page_content="\n".join(linhas), metadata={"source": os.path.basename(caminho_arquivo), "format": "HTML (.html)"})]

def extrair_documentos_markdown(caminho_arquivo: str) -> List[Document]:
    with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
        conteudo = f.read()
    return [Document(page_content=conteudo, metadata={"source": os.path.basename(caminho_arquivo), "format": "Markdown (.md)"})]

def carregar_arquivo_por_extensao(caminho_arquivo: str) -> List[Document]:
    """
    Função unificada que detecta a extensão e aplica o carregador apropriado
    para os 8 formatos suportados pelo Challenge Alura Agentes.
    """
    extensao = os.path.splitext(caminho_arquivo)[1].lower()
    
    try:
        if extensao == '.pdf':
            return extrair_documentos_pdf(caminho_arquivo)
        elif extensao in ['.docx', '.doc']:
            return extrair_documentos_docx(caminho_arquivo)
        elif extensao in ['.xlsx', '.xls']:
            return extrair_documentos_excel(caminho_arquivo)
        elif extensao in ['.pptx', '.ppt']:
            return extrair_documentos_pptx(caminho_arquivo)
        elif extensao == '.csv':
            return extrair_documentos_csv(caminho_arquivo)
        elif extensao == '.json':
            return extrair_documentos_json(caminho_arquivo)
        elif extensao in ['.html', '.htm']:
            return extrair_documentos_html(caminho_arquivo)
        elif extensao in ['.md', '.txt']:
            return extrair_documentos_markdown(caminho_arquivo)
        else:
            # Fallback para arquivos de texto simples
            return extrair_documentos_markdown(caminho_arquivo)
    except Exception as e:
        print(f"Erro ao processar o arquivo {caminho_arquivo}: {e}")
        return []
